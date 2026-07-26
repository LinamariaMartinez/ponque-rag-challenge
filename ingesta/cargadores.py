"""Carga el texto plano de un documento según su formato."""

import csv
import json
from pathlib import Path

from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader


def cargar_markdown(ruta: Path) -> str:
    return ruta.read_text(encoding="utf-8")


def cargar_csv(ruta: Path) -> str:
    with ruta.open(encoding="utf-8") as archivo:
        lector = csv.reader(archivo)
        return "\n".join(" | ".join(fila) for fila in lector)


def cargar_json(ruta: Path) -> str:
    datos = json.loads(ruta.read_text(encoding="utf-8"))
    return json.dumps(datos, ensure_ascii=False, indent=2)


def cargar_html(ruta: Path) -> str:
    sopa = BeautifulSoup(ruta.read_text(encoding="utf-8"), "html.parser")
    return sopa.get_text(separator="\n", strip=True)


def cargar_pdf(ruta: Path) -> str:
    lector = PdfReader(str(ruta))
    return "\n".join(pagina.extract_text() or "" for pagina in lector.pages)


def cargar_docx(ruta: Path) -> str:
    documento = DocxDocument(str(ruta))
    return "\n".join(p.text for p in documento.paragraphs if p.text.strip())


def cargar_xlsx(ruta: Path) -> str:
    libro = load_workbook(str(ruta), data_only=True)
    lineas = []
    for hoja in libro.worksheets:
        lineas.append(f"Hoja: {hoja.title}")
        for fila in hoja.iter_rows(values_only=True):
            valores = [str(v) for v in fila if v is not None]
            if valores:
                lineas.append(" | ".join(valores))
    return "\n".join(lineas)


def cargar_pptx(ruta: Path) -> str:
    presentacion = Presentation(str(ruta))
    lineas = []
    for i, diapositiva in enumerate(presentacion.slides, start=1):
        lineas.append(f"Diapositiva {i}")
        for forma in diapositiva.shapes:
            if forma.has_text_frame and forma.text_frame.text.strip():
                lineas.append(forma.text_frame.text.strip())
    return "\n".join(lineas)


CARGADORES = {
    ".md": cargar_markdown,
    ".csv": cargar_csv,
    ".json": cargar_json,
    ".html": cargar_html,
    ".pdf": cargar_pdf,
    ".docx": cargar_docx,
    ".xlsx": cargar_xlsx,
    ".pptx": cargar_pptx,
}


def cargar_documento(ruta: Path) -> str:
    cargador = CARGADORES.get(ruta.suffix.lower())
    if cargador is None:
        raise ValueError(f"Formato no soportado: {ruta.suffix}")
    return cargador(ruta)
