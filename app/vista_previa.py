"""Vista previa formateada de un documento (para el panel lateral, no para el RAG).

A diferencia de ingesta/cargadores.py (que aplana todo a texto plano para
embeber), aquí se preserva la estructura de cada formato: tablas reales,
HTML renderizado, o una lista de párrafos."""

import csv
import json
from pathlib import Path

import markdown
from bs4 import BeautifulSoup
from docx import Document as DocxDocument
from openpyxl import load_workbook
from pptx import Presentation
from pypdf import PdfReader

from app.config import DOCUMENTOS_DIR


def _vista_previa_csv(ruta: Path) -> dict:
    with ruta.open(encoding="utf-8") as archivo:
        filas = list(csv.reader(archivo))
    encabezados = filas[0] if filas else []
    resto = filas[1:] if len(filas) > 1 else []
    return {"tipo": "tabla", "contenido": [{"hoja": None, "encabezados": encabezados, "filas": resto}]}


def _vista_previa_xlsx(ruta: Path) -> dict:
    libro = load_workbook(str(ruta), data_only=True)
    hojas = []
    for hoja in libro.worksheets:
        filas_todas = [
            [("" if v is None else str(v)) for v in fila]
            for fila in hoja.iter_rows(values_only=True)
        ]
        encabezados = filas_todas[0] if filas_todas else []
        resto = filas_todas[1:] if len(filas_todas) > 1 else []
        hojas.append({"hoja": hoja.title, "encabezados": encabezados, "filas": resto})
    return {"tipo": "tabla", "contenido": hojas}


def _vista_previa_markdown(ruta: Path) -> dict:
    texto = ruta.read_text(encoding="utf-8")
    return {"tipo": "html", "contenido": markdown.markdown(texto)}


def _vista_previa_json(ruta: Path) -> dict:
    datos = json.loads(ruta.read_text(encoding="utf-8"))
    return {"tipo": "texto", "contenido": [json.dumps(datos, ensure_ascii=False, indent=2)]}


def _vista_previa_docx(ruta: Path) -> dict:
    documento = DocxDocument(str(ruta))
    parrafos = [p.text for p in documento.paragraphs if p.text.strip()]
    return {"tipo": "texto", "contenido": parrafos}


def _vista_previa_pdf(ruta: Path) -> dict:
    lector = PdfReader(str(ruta))
    paginas = [pagina.extract_text() or "" for pagina in lector.pages]
    paginas = [p for p in paginas if p.strip()]
    return {"tipo": "texto", "contenido": paginas}


def _vista_previa_pptx(ruta: Path) -> dict:
    presentacion = Presentation(str(ruta))
    bloques = []
    for i, diapositiva in enumerate(presentacion.slides, start=1):
        lineas = [f"Diapositiva {i}"]
        for forma in diapositiva.shapes:
            if forma.has_text_frame and forma.text_frame.text.strip():
                lineas.append(forma.text_frame.text.strip())
        bloques.append("\n".join(lineas))
    return {"tipo": "texto", "contenido": bloques}


def _vista_previa_html(ruta: Path) -> dict:
    sopa = BeautifulSoup(ruta.read_text(encoding="utf-8"), "html.parser")
    cuerpo = sopa.find("body")
    contenido = cuerpo.decode_contents() if cuerpo else str(sopa)
    return {"tipo": "html", "contenido": contenido}


_GENERADORES_VISTA_PREVIA = {
    ".csv": _vista_previa_csv,
    ".xlsx": _vista_previa_xlsx,
    ".md": _vista_previa_markdown,
    ".json": _vista_previa_json,
    ".docx": _vista_previa_docx,
    ".pdf": _vista_previa_pdf,
    ".pptx": _vista_previa_pptx,
    ".html": _vista_previa_html,
}


def generar_vista_previa(area: str, archivo: str) -> dict:
    ruta = (DOCUMENTOS_DIR / area / archivo).resolve()
    if DOCUMENTOS_DIR.resolve() not in ruta.parents or not ruta.is_file():
        raise FileNotFoundError(f"No existe: {area}/{archivo}")
    generador = _GENERADORES_VISTA_PREVIA.get(ruta.suffix.lower())
    if generador is None:
        raise ValueError(f"Formato no soportado: {ruta.suffix}")
    return generador(ruta)
