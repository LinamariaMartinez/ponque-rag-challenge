import tempfile
import unittest
from pathlib import Path

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from ingesta.cargadores import cargar_documento


class TestCargadoresBinarios(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)

    def tearDown(self):
        self.tmp.cleanup()

    def test_carga_pdf(self):
        ruta = self.dir / "doc.pdf"
        lienzo = canvas.Canvas(str(ruta), pagesize=letter)
        lienzo.drawString(72, 720, "Contenido de prueba PDF")
        lienzo.save()
        self.assertIn("Contenido de prueba PDF", cargar_documento(ruta))

    def test_carga_docx(self):
        ruta = self.dir / "doc.docx"
        documento = Document()
        documento.add_paragraph("Contenido de prueba DOCX")
        documento.save(str(ruta))
        self.assertIn("Contenido de prueba DOCX", cargar_documento(ruta))

    def test_carga_xlsx(self):
        ruta = self.dir / "doc.xlsx"
        libro = Workbook()
        hoja = libro.active
        hoja.append(["Categoría", "Valor"])
        hoja.append(["Prueba", 100])
        libro.save(str(ruta))
        texto = cargar_documento(ruta)
        self.assertIn("Prueba", texto)
        self.assertIn("100", texto)

    def test_carga_pptx(self):
        ruta = self.dir / "doc.pptx"
        presentacion = Presentation()
        diapositiva = presentacion.slides.add_slide(presentacion.slide_layouts[1])
        diapositiva.shapes.title.text = "Contenido de prueba PPTX"
        presentacion.save(str(ruta))
        self.assertIn("Contenido de prueba PPTX", cargar_documento(ruta))
