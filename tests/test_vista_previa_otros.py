import tempfile
import unittest
from pathlib import Path
from unittest.mock import patch

from docx import Document
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from app.vista_previa import generar_vista_previa


class TestVistaPreviaOtros(unittest.TestCase):
    def setUp(self):
        self.tmp = tempfile.TemporaryDirectory()
        self.dir = Path(self.tmp.name)
        (self.dir / "rh").mkdir()
        self.parche = patch("app.vista_previa.DOCUMENTOS_DIR", self.dir)
        self.parche.start()

    def tearDown(self):
        self.parche.stop()
        self.tmp.cleanup()

    def test_vista_previa_markdown(self):
        ruta = self.dir / "rh" / "notas.md"
        ruta.write_text("# Título\n\nContenido de prueba.", encoding="utf-8")
        resultado = generar_vista_previa("rh", "notas.md")
        self.assertEqual(resultado["tipo"], "html")
        self.assertIn("<h1>Título</h1>", resultado["contenido"])
        self.assertIn("Contenido de prueba", resultado["contenido"])

    def test_vista_previa_json(self):
        ruta = self.dir / "rh" / "datos.json"
        ruta.write_text('{"tema": "prueba"}', encoding="utf-8")
        resultado = generar_vista_previa("rh", "datos.json")
        self.assertEqual(resultado["tipo"], "texto")
        self.assertEqual(len(resultado["contenido"]), 1)
        self.assertIn('"tema": "prueba"', resultado["contenido"][0])

    def test_vista_previa_docx(self):
        ruta = self.dir / "rh" / "politica.docx"
        documento = Document()
        documento.add_paragraph("Primer párrafo.")
        documento.add_paragraph("Segundo párrafo.")
        documento.save(str(ruta))
        resultado = generar_vista_previa("rh", "politica.docx")
        self.assertEqual(resultado["tipo"], "texto")
        self.assertEqual(resultado["contenido"], ["Primer párrafo.", "Segundo párrafo."])

    def test_vista_previa_pdf(self):
        ruta = self.dir / "rh" / "manual.pdf"
        lienzo = canvas.Canvas(str(ruta), pagesize=letter)
        lienzo.drawString(72, 720, "Contenido de la página")
        lienzo.save()
        resultado = generar_vista_previa("rh", "manual.pdf")
        self.assertEqual(resultado["tipo"], "texto")
        self.assertIn("Contenido de la página", resultado["contenido"][0])

    def test_vista_previa_pptx(self):
        ruta = self.dir / "rh" / "plan.pptx"
        presentacion = Presentation()
        diapositiva = presentacion.slides.add_slide(presentacion.slide_layouts[1])
        diapositiva.shapes.title.text = "Título de prueba"
        presentacion.save(str(ruta))
        resultado = generar_vista_previa("rh", "plan.pptx")
        self.assertEqual(resultado["tipo"], "texto")
        self.assertIn("Diapositiva 1", resultado["contenido"][0])
        self.assertIn("Título de prueba", resultado["contenido"][0])

    def test_vista_previa_html(self):
        ruta = self.dir / "rh" / "boletin.html"
        ruta.write_text(
            "<html><body><h1>Aviso</h1><p>Texto de prueba</p></body></html>",
            encoding="utf-8",
        )
        resultado = generar_vista_previa("rh", "boletin.html")
        self.assertEqual(resultado["tipo"], "html")
        self.assertIn("Aviso", resultado["contenido"])
        self.assertIn("Texto de prueba", resultado["contenido"])
