"""Genera los documentos ficticios en formato binario (DOCX, XLSX, PPTX, PDF)
para documentos/. Los formatos de texto plano (JSON, CSV, MD, HTML) ya están
creados directamente como archivos y no pasan por este script."""

import sys
from pathlib import Path

sys.path.insert(0, str(Path(__file__).resolve().parent.parent))

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from reportlab.lib.pagesizes import letter
from reportlab.pdfgen import canvas

from app.config import DOCUMENTOS_DIR


def generar_politica_vacaciones() -> None:
    ruta = DOCUMENTOS_DIR / "rh" / "politica-vacaciones.docx"
    ruta.parent.mkdir(parents=True, exist_ok=True)
    documento = Document()
    documento.add_heading("Política de Vacaciones — Ponqué Ponqué Calarcá", level=1)
    documento.add_paragraph("Documento de ejemplo (datos ficticios) para el agente RAG.")
    documento.add_heading("Días disponibles", level=2)
    documento.add_paragraph(
        "Cada colaborador con más de un año de antigüedad tiene derecho a 15 días "
        "hábiles de vacaciones remuneradas por año."
    )
    documento.add_heading("Cómo solicitarlas", level=2)
    documento.add_paragraph(
        "La solicitud se hace por escrito al área de Recursos Humanos con al menos "
        "15 días calendario de anticipación. RH confirma la aprobación en un plazo "
        "máximo de 3 días hábiles."
    )
    documento.add_heading("Pago", level=2)
    documento.add_paragraph(
        "El pago de vacaciones se realiza junto con la nómina del mes en que inicia "
        "el periodo de descanso."
    )
    documento.save(str(ruta))
    print(f"  [ok] {ruta.relative_to(DOCUMENTOS_DIR.parent)}")


def generar_presupuesto() -> None:
    ruta = DOCUMENTOS_DIR / "finanzas" / "presupuesto-trimestral.xlsx"
    ruta.parent.mkdir(parents=True, exist_ok=True)
    libro = Workbook()
    hoja = libro.active
    hoja.title = "Presupuesto Q3 2026"
    hoja.append(["Categoría", "Presupuestado (COP)", "Ejecutado (COP)"])
    filas = [
        ("Materia prima fábrica", 18000000, 17250000),
        ("Transporte y logística", 4500000, 4800000),
        ("Publicidad", 3000000, 2100000),
        ("Nómina", 12000000, 12000000),
        ("Empaques", 2200000, 2350000),
    ]
    for fila in filas:
        hoja.append(list(fila))
    libro.save(str(ruta))
    print(f"  [ok] {ruta.relative_to(DOCUMENTOS_DIR.parent)}")


def generar_plan_campana() -> None:
    ruta = DOCUMENTOS_DIR / "marketing" / "plan-campana-dia-madre.pptx"
    ruta.parent.mkdir(parents=True, exist_ok=True)
    presentacion = Presentation()

    diapositiva = presentacion.slides.add_slide(presentacion.slide_layouts[0])
    diapositiva.shapes.title.text = "Campaña Día de la Madre 2026 (ejemplo)"
    diapositiva.placeholders[1].text = "Ponqué Ponqué Calarcá — plan de marketing ficticio"

    layout_contenido = presentacion.slide_layouts[1]

    diapositiva = presentacion.slides.add_slide(layout_contenido)
    diapositiva.shapes.title.text = "Objetivo"
    diapositiva.placeholders[1].text = (
        "Aumentar en 20% los pedidos de la semana del Día de la Madre frente al mes "
        "anterior (meta ficticia de ejemplo)."
    )

    diapositiva = presentacion.slides.add_slide(layout_contenido)
    diapositiva.shapes.title.text = "Canales"
    marco = diapositiva.placeholders[1].text_frame
    marco.text = "Instagram (pauta + orgánico)"
    marco.add_paragraph().text = "WhatsApp (catálogo a clientes frecuentes)"
    marco.add_paragraph().text = "Volanteo en puntos de venta aliados"

    diapositiva = presentacion.slides.add_slide(layout_contenido)
    diapositiva.shapes.title.text = "Presupuesto y cronograma"
    marco = diapositiva.placeholders[1].text_frame
    marco.text = "Presupuesto: $1.500.000 COP (ficticio)"
    marco.add_paragraph().text = "Del 20 de abril al 10 de mayo de 2026"

    presentacion.save(str(ruta))
    print(f"  [ok] {ruta.relative_to(DOCUMENTOS_DIR.parent)}")


def generar_manual_despacho() -> None:
    ruta = DOCUMENTOS_DIR / "operaciones" / "manual-despacho.pdf"
    ruta.parent.mkdir(parents=True, exist_ok=True)
    lienzo = canvas.Canvas(str(ruta), pagesize=letter)
    _, alto = letter
    y = alto - 72

    def linea(texto: str, salto: int = 18, tamano: int = 11) -> None:
        nonlocal y
        lienzo.setFont("Helvetica", tamano)
        lienzo.drawString(72, y, texto)
        y -= salto

    linea("Manual de Empaque y Despacho — Ponqué Ponqué Calarcá", tamano=14)
    linea("Documento de ejemplo (datos ficticios) para el agente RAG.", salto=28)
    linea("1. Verificar el pedido contra la orden antes de empacar.")
    linea("2. Usar caja rígida para transportes de más de 30 minutos.")
    linea("3. Sellar la caja con cinta de seguridad con el logo de la marca.")
    linea("4. Confirmar la dirección y el nombre del cliente antes de despachar.")
    linea("5. Registrar la hora de salida del domiciliario en la planilla diaria.")
    lienzo.save()
    print(f"  [ok] {ruta.relative_to(DOCUMENTOS_DIR.parent)}")


def main() -> None:
    print("Generando documentos ficticios binarios...")
    generar_politica_vacaciones()
    generar_presupuesto()
    generar_plan_campana()
    generar_manual_despacho()
    print("Listo.")


if __name__ == "__main__":
    main()
